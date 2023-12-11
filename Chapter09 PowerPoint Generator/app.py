import collections.abc
import config

assert collections
import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import openai
from io import BytesIO
import requests

# API Token
openai.api_key = config.API_KEY


def slide_generator(text, prs):
    prompt = f"Summarize the following text to a DALL-E image generation suitable as a background image in slides" \
             f"prompt: \n {text}"

    model_engine = "gpt-4"
    dlp = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )
    print(dlp)

    dalle_prompt = dlp["choices"][0]["message"]["content"]

    response = openai.Image.create(
        prompt=dalle_prompt + " Style: digital art",
        n=1,
        size="1024x1024"
    )
    print(response)
    image_url = response['data'][0]['url']

    prompt = f"Create a bullet point text for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    print(ppt)
    ppt_text = ppt["choices"][0]["message"]["content"]

    prompt = f"Create a title for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = openai.ChatCompletion.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    print(ppt)
    ppt_header = ppt["choices"][0]["message"]["content"]

    # Add a new slide to the presentation
    slide = prs.slides.add_slide(prs.slide_layouts[3])

    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(0), Inches(0))

    # Add text box
    txBox = slide.shapes.add_textbox(Inches(3), Inches(1),
                                     Inches(4), Inches(1.5))
    tf = txBox.text_frame
    tf.text = ppt_text

    # title
    slide.shapes.title.text = ppt_header
    # subtitle
    slide.placeholders[1].text = "Quite Cool"


def get_slides():
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height
    print("# of paragraphs: ", len(paragraphs))
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)

    prs.save("my_presentation.pptx")
    print("PPT saved!")


def add_background_image():
    pass


def add_table(prs):
    title_only_slide_layout=prs.slide_layouts[5]
    #根据上面版式新建一张PPT
    slide=prs.slides.add_slide(title_only_slide_layout)
    #把新建PPT中所有的形状集合赋值给变量shapes
    shapes=slide.shapes
    #设置标题的文字
    shapes.title.text='Adding a Table'
    
    #设计将要新建的表格的行列数
    rows=cols=2
    #设置距离幻灯片左边及顶端的距离
    left=top=Inches(2.0)
    #设置表格的宽度
    width=Inches(6.0)
    #设置表格的高度
    height=Inches(0.8)
    #根据指定的行列数、位置、大小新建一个表格
    table=shapes.add_table(rows,cols,left,top,width,height).table
    #重新调整每列的宽度
    table.columns[0].width=Inches(2.0)
    table.columns[1].widht=Inches(4.0)
    
    #设置表头的文字
    table.cell(0,0).text='Foo'
    table.cell(0,1).text='Bar'
    
    #设置表体单元格里的文字
    table.cell(1,0).text='Baz'
    table.cell(1,1).text='Qux'


def add_shapes(prs):
    title_only_slide_layout=prs.slide_layouts[5]
    #根据指定样式新建一张PPT，
    slide=prs.slides.add_slide(title_only_slide_layout)
    #把新建PPT中的所有形状对象shapes赋值给变量shapes
    shapes=slide.shapes
    #设置标题的文字
    shapes.title.text='Adding an AutoShape'
    #设置下面新建形状距所在页PPT左边的位置
    left=Inches(0.93)
    #设置下面新建形状距所在页PPT顶商的位置
    top=Inches(3.0)
    #设置插入形状的宽度
    width=Inches(1.75)
    #设置插入形状的高度
    height=Inches(1.0)
    #根据指定位置、大小新建一个五角形
    shape=shapes.add_shape(MSO_SHAPE.PENTAGON,left,top,width,height)
    #给新建的形状添加文字
    shape.text='Step 1'
    #设置下面新建V形距左边的位置，0.4是V形的深度
    left=left+width-Inches(0.4)
    #设置新建V形的宽度
    width=Inches(2.0)
    #通过循环创建4个V形形状
    for n in range(2,6):
        #根据指定的位置、大小新建一个形状
        shape=shapes.add_shape(MSO_SHAPE.CHEVRON,left,top,width,height)
        #形状的文字
        shape.text='Step %d'%n
        #重新调整下一个形状距离左边的位置
        left=left+width-Inches(0.4)


if __name__ == "__main__":

    app = tk.Tk()
    app.title("Crate PPT Slides")
    app.geometry("800x600")

    # Create text field
    text_field = tk.Text(app)
    text_field.pack(fill="both", expand=True)
    text_field.configure(wrap="word", font=("Arial", 12))
    text_field.focus_set()

    # Create the button to create slides
    create_button = tk.Button(app, text="Create Slides", command=get_slides)
    create_button.pack()

    app.mainloop()
